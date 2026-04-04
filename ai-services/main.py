from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
import os
import tempfile
import boto3
from botocore.config import Config

app = FastAPI(title="Brainwave.ai AI Services", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


def get_s3():
    return boto3.client(
        "s3",
        endpoint_url=os.getenv("MINIO_ENDPOINT_URL", "http://minio:9000"),
        aws_access_key_id=os.getenv("MINIO_ACCESS_KEY", "minioadmin"),
        aws_secret_access_key=os.getenv("MINIO_SECRET_KEY", "minioadmin123"),
        config=Config(signature_version="s3v4"),
        region_name="us-east-1"
    )


class TranscribeRequest(BaseModel):
    file_url: str
    material_id: str
    language: Optional[str] = None  # Whisper language code or None for auto


class StructureRequest(BaseModel):
    content: str
    language: Optional[str] = None


class QuizRequest(BaseModel):
    content: str
    num_questions: int = 5
    language: Optional[str] = None


class ModerateRequest(BaseModel):
    title: str
    category: str = ""
    body_text: str = ""
    file_names: Optional[list] = None


class ThumbnailRequest(BaseModel):
    title: str
    category: str
    description: str = ""
    lesson_title: Optional[str] = None
    module_title: Optional[str] = None
    faculty_face_image_url: Optional[str] = None
    custom_prompt: Optional[str] = None


class EmbedRequest(BaseModel):
    text: str


class ChatRequest(BaseModel):
    course_id: str
    course_name: str
    question: str


class IndexRequest(BaseModel):
    course_id: str
    text: str
    lesson_id: str
    chapter_id: str
    source_type: str = "course_content"


class RiskRequest(BaseModel):
    avg_watch_percent: float = 0
    rewatch_rate: float = 0
    quiz_avg_score: float = 0
    days_since_last_active: int = 0
    completion_rate: float = 0
    community_posts_count: int = 0
    assignment_submit_rate: float = 0
    enrollment_days: int = 0


class ExtractTextRequest(BaseModel):
    file_url: str
    file_type: str


@app.get("/health")
async def health():
    return {"status": "ok", "service": "AI Services"}


@app.post("/transcribe")
async def transcribe(data: TranscribeRequest):
    """Transcribe audio/video from URL."""
    from transcription import transcribe_from_url
    try:
        result = transcribe_from_url(data.file_url, data.material_id, language=data.language)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/extract-text")
async def extract_text(data: ExtractTextRequest):
    """Extract text from PDF/DOC/PPT."""
    import httpx

    with tempfile.NamedTemporaryFile(suffix=f".{data.file_type}", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(data.file_url)
            with open(tmp_path, "wb") as f:
                f.write(resp.content)

        text = ""
        if data.file_type == "pdf":
            import PyPDF2
            with open(tmp_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""

        elif data.file_type in ["doc", "docx"]:
            from docx import Document
            doc = Document(tmp_path)
            text = "\n".join([p.text for p in doc.paragraphs])

        elif data.file_type in ["ppt", "pptx"]:
            from pptx import Presentation
            prs = Presentation(tmp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        return {"text": text}
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


@app.post("/moderate-content")
async def moderate_content(data: ModerateRequest):
    from moderate_content import moderate_course_content
    result = await moderate_course_content(
        data.title, data.category, data.body_text, data.file_names or []
    )
    return result


@app.post("/structure-course")
async def structure_course(data: StructureRequest):
    """Structure transcript into course chapters using Llama 3."""
    from course_builder import structure_course as _structure
    result = await _structure(data.content, language=data.language)
    return result


@app.post("/generate-quiz")
async def generate_quiz(data: QuizRequest):
    """Generate quiz questions from content using Llama 3."""
    from quiz_generator import generate_quiz as _generate
    result = await _generate(data.content, data.num_questions, language=data.language)
    return result


@app.post("/embed")
async def embed(data: EmbedRequest):
    """Get text embedding vector."""
    from embeddings import embed_text
    vector = await embed_text(data.text)
    return {"embedding": vector, "dimensions": len(vector)}


@app.post("/chat")
async def chat(data: ChatRequest):
    """RAG chatbot for course-specific Q&A."""
    from chatbot import chat as _chat
    result = await _chat(data.course_id, data.course_name, data.question)
    return result


@app.post("/index")
async def index_content(data: IndexRequest):
    """Index course material into Qdrant."""
    from indexer import index_course_material
    await index_course_material(
        course_id=data.course_id,
        text=data.text,
        lesson_id=data.lesson_id,
        chapter_id=data.chapter_id,
        source_type=data.source_type
    )
    return {"message": "Content indexed successfully"}


@app.post("/generate-thumbnail")
async def generate_thumbnail(data: ThumbnailRequest):
    """Generate course thumbnail using Gemini or PIL fallback."""
    from thumbnail_generator import generate_thumbnail as _generate
    from app_storage import upload_thumbnail

    try:
        img_bytes = await _generate(
            data.title,
            data.category,
            data.description,
            lesson_title=data.lesson_title,
            module_title=data.module_title,
            faculty_face_image_url=data.faculty_face_image_url,
            custom_prompt=data.custom_prompt,
        )
        # Upload to MinIO
        import uuid
        key = f"ai_generated/{uuid.uuid4()}.png"
        url = upload_thumbnail(img_bytes, key)
        return {"thumbnail_url": url}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/predict-dropout")
async def predict_dropout(data: RiskRequest):
    """Predict student dropout risk."""
    from dropout_predictor import predict_risk
    features = data.dict()
    result = predict_risk(features)
    return result


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8001, reload=True)
